"""Tests for PPTX whole-slide screenshot rendering, notes extraction, and VLM analysis."""

from __future__ import annotations

import threading
import time
from pathlib import Path
from unittest.mock import patch

import pytest
import typer
from PIL import Image
from pptx import Presentation

from config import Settings
from doc_convert.pptx_slide_vlm import (
    DEFAULT_PPTX_SLIDE_VLM,
    _normalize_heading_levels,
    _pptx_to_pdf,
    analyze_slide_images,
    extract_hidden_slide_numbers,
    extract_slide_notes,
)


def _make_pptx(tmp_path: Path, *, notes: dict[int, str]) -> Path:
    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank layout
    slide_count = max(notes.keys(), default=1)
    for i in range(1, slide_count + 1):
        slide = prs.slides.add_slide(layout)
        if i in notes:
            slide.notes_slide.notes_text_frame.text = notes[i]
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def test_default_model_is_confirmed_ibm_sonnet() -> None:
    # Confirmed present via IBM ICA's GET /models (see providers.DEFAULT_PPTX_SLIDE_VLM).
    assert DEFAULT_PPTX_SLIDE_VLM == "ibm/claude-sonnet-4-6"


def test_extract_slide_notes_returns_only_nonempty(tmp_path: Path) -> None:
    pptx_path = _make_pptx(tmp_path, notes={1: "Remember to mention the roadmap", 3: "  "})
    notes = extract_slide_notes(pptx_path)
    assert notes == {1: "Remember to mention the roadmap"}


def test_extract_hidden_slide_numbers_marks_hidden_only(tmp_path: Path) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[6]
    for _ in range(3):
        prs.slides.add_slide(layout)
    # python-pptx has no high-level API for the hidden flag; set the raw XML
    # attribute the same way PowerPoint does (<p:sld show="0">).
    prs.slides[1]._element.set("show", "0")
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    assert extract_hidden_slide_numbers(path) == {2}


def test_extract_hidden_slide_numbers_none_hidden(tmp_path: Path) -> None:
    pptx_path = _make_pptx(tmp_path, notes={})
    assert extract_hidden_slide_numbers(pptx_path) == set()


def test_pptx_to_pdf_forces_export_hidden_slides(tmp_path: Path) -> None:
    """LibreOffice must be told to keep hidden slides in the PDF: otherwise PDF
    page numbers drift out of sync with docling/python-pptx slide numbering,
    misaligning every subsequent slide's text/visual/notes sections.
    """
    pptx_path = tmp_path / "deck.pptx"
    pptx_path.write_bytes(b"fake")
    captured_cmd: list[str] = []

    class _Result:
        returncode = 0
        stderr = ""

    def fake_run(cmd: list[str], **_kwargs: object) -> _Result:
        captured_cmd.extend(cmd)
        (tmp_path / f"{pptx_path.stem}.pdf").write_bytes(b"%PDF-1.4 fake")
        return _Result()

    with patch("shutil.which", return_value="/usr/bin/soffice"), patch("subprocess.run", side_effect=fake_run):
        _pptx_to_pdf(pptx_path, tmp_path)

    convert_arg = captured_cmd[captured_cmd.index("--convert-to") + 1]
    assert "ExportHiddenSlides" in convert_arg
    assert '"value":"true"' in convert_arg


def test_extract_slide_notes_no_notes_slide(tmp_path: Path) -> None:
    pptx_path = _make_pptx(tmp_path, notes={})
    assert extract_slide_notes(pptx_path) == {}


def test_pptx_to_pdf_missing_soffice_exits(tmp_path: Path) -> None:
    with patch("shutil.which", return_value=None), pytest.raises(typer.Exit):
        _pptx_to_pdf(tmp_path / "deck.pptx", tmp_path)


# ── heading-level safety net (VLM sometimes emits its own top-level heading) ───


def test_normalize_heading_levels_demotes_h2() -> None:
    text = "## Slide Analysis\n\nSome content."
    out = _normalize_heading_levels(text)
    assert out.startswith("#### Slide Analysis")
    assert "## Slide Analysis" not in out.replace("#### Slide Analysis", "")


def test_normalize_heading_levels_leaves_deep_headings_alone() -> None:
    text = "#### Purpose\n\nContent\n\n##### Sub-point"
    assert _normalize_heading_levels(text) == text


def test_normalize_heading_levels_no_heading_is_noop() -> None:
    text = "Just plain text, no heading markers at all."
    assert _normalize_heading_levels(text) == text


def test_normalize_heading_levels_multiple_offending_lines() -> None:
    text = "# Title\n\nIntro.\n\n## Section\n\nBody."
    out = _normalize_heading_levels(text)
    assert out.startswith("#### Title\n")
    assert "#### Section" in out
    assert not out.splitlines()[0].startswith("# ")
    assert "\n## Section" not in out


# ---------------------------------------------------------------------------
# Concurrent slide analysis
# ---------------------------------------------------------------------------


def _slide_files(tmp_path: Path, count: int) -> dict[int, str]:
    """Create ``count`` fake slide PNGs and return the {slide_no: rel_path} map."""
    slides = tmp_path / "slides"
    slides.mkdir(parents=True, exist_ok=True)
    mapping: dict[int, str] = {}
    for i in range(1, count + 1):
        name = f"slide_{i:03d}.png"
        Image.new("RGB", (80, 60), (i * 10 % 255, 40, 90)).save(slides / name)
        mapping[i] = f"slides/{name}"
    return mapping


def test_analyze_slide_images_runs_concurrently_and_keys_by_slide(tmp_path: Path, ibm_settings: Settings) -> None:
    """All slides are analysed and mapped to the right number, whatever the
    completion order (which is nondeterministic under concurrency)."""
    mapping = _slide_files(tmp_path, 9)
    seen: list[int] = []
    max_in_flight = 0
    in_flight = 0
    lock = threading.Lock()

    def fake_describe(image_path: Path, *_args: object, **_kwargs: object) -> str:
        nonlocal in_flight, max_in_flight
        with lock:
            in_flight += 1
            max_in_flight = max(max_in_flight, in_flight)
            seen.append(int(image_path.stem.split("_")[1]))
        time.sleep(0.05)  # hold the slot so overlap is observable
        with lock:
            in_flight -= 1
        return f"analysis of {image_path.name}"

    with patch("doc_convert.vision_llm.describe_image", side_effect=fake_describe):
        analyses = analyze_slide_images(mapping, tmp_path, ibm_settings, concurrency=4)

    assert sorted(analyses) == list(range(1, 10))
    assert sorted(seen) == list(range(1, 10))
    assert max_in_flight > 1, "requests must actually overlap"
    assert max_in_flight <= 4, "must never exceed the requested concurrency"
    # Each slide keeps its own description: no cross-contamination between threads.
    for number, analysis in analyses.items():
        assert f"slide_{number:03d}.png" in analysis.description
        assert analysis.slide_number == number
        assert analysis.image_path == mapping[number]


def test_analyze_slide_images_concurrency_one_is_sequential(tmp_path: Path, ibm_settings: Settings) -> None:
    """concurrency=1 must never overlap requests (escape hatch for rate limits)."""
    mapping = _slide_files(tmp_path, 5)
    max_in_flight = 0
    in_flight = 0
    lock = threading.Lock()

    def fake_describe(image_path: Path, *_args: object, **_kwargs: object) -> str:
        nonlocal in_flight, max_in_flight
        with lock:
            in_flight += 1
            max_in_flight = max(max_in_flight, in_flight)
        time.sleep(0.02)
        with lock:
            in_flight -= 1
        return "ok"

    with patch("doc_convert.vision_llm.describe_image", side_effect=fake_describe):
        analyses = analyze_slide_images(mapping, tmp_path, ibm_settings, concurrency=1)

    assert max_in_flight == 1
    assert len(analyses) == 5


def test_analyze_slide_images_result_is_ordered_by_slide_number(tmp_path: Path, ibm_settings: Settings) -> None:
    """Completion order is arbitrary; the returned mapping must still iterate in
    slide order so logs and debugging stay readable."""
    mapping = _slide_files(tmp_path, 6)

    def fake_describe(image_path: Path, *_args: object, **_kwargs: object) -> str:
        # Make later slides finish first.
        time.sleep(max(0.0, (7 - int(image_path.stem.split("_")[1])) * 0.01))
        return "ok"

    with patch("doc_convert.vision_llm.describe_image", side_effect=fake_describe):
        analyses = analyze_slide_images(mapping, tmp_path, ibm_settings, concurrency=4)

    assert list(analyses) == sorted(analyses)


def test_analyze_slide_images_one_failure_does_not_lose_the_others(tmp_path: Path, ibm_settings: Settings) -> None:
    """A slide that exhausts its retries is simply absent; the deck still converts."""
    mapping = _slide_files(tmp_path, 4)

    def fake_describe(image_path: Path, *_args: object, **_kwargs: object) -> str:
        if image_path.stem.endswith("002"):
            return ""  # exhausted its retries inside the primitive
        return "ok"

    with patch("doc_convert.vision_llm.describe_image", side_effect=fake_describe):
        analyses = analyze_slide_images(mapping, tmp_path, ibm_settings, concurrency=4)

    assert sorted(analyses) == [1, 3, 4]


def test_analyze_slide_images_skips_missing_screenshots(tmp_path: Path, ibm_settings: Settings) -> None:
    mapping = _slide_files(tmp_path, 2)
    mapping[3] = "slides/slide_003.png"  # never written to disk

    def fake_describe(image_path: Path, *_args: object, **_kwargs: object) -> str:
        assert image_path.exists(), "a missing file must never reach the API call"
        return "ok"

    with patch("doc_convert.vision_llm.describe_image", side_effect=fake_describe):
        analyses = analyze_slide_images(mapping, tmp_path, ibm_settings, concurrency=4)

    assert sorted(analyses) == [1, 2]


def test_analyze_slide_images_no_slides_returns_empty(tmp_path: Path, ibm_settings: Settings) -> None:
    assert analyze_slide_images({}, tmp_path, ibm_settings) == {}


def test_analyze_slide_images_fatal_model_error_aborts(tmp_path: Path, ibm_settings: Settings) -> None:
    """An unknown model slug must abort the whole deck rather than silently
    producing 52 empty analyses."""
    mapping = _slide_files(tmp_path, 8)

    def fake_describe(_image_path: Path, *_args: object, **_kwargs: object) -> str:
        raise typer.Exit(1)

    with (
        patch("doc_convert.vision_llm.describe_image", side_effect=fake_describe),
        pytest.raises(typer.Exit),
    ):
        analyze_slide_images(mapping, tmp_path, ibm_settings, concurrency=4)
