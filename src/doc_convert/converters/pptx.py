"""PPTX converter: Docling native + python-pptx images + VLM descriptions."""

from __future__ import annotations

import hashlib
import logging
from pathlib import Path  # noqa: TC003

from doc_convert.base import BaseConverter
from doc_convert.markdown import (
    FloatingArtifacts,
    FloatingContext,
    build_document_markdown,
    build_images_catalog,
    collect_floating_contexts,
    extract_title,
)

logger = logging.getLogger(__name__)

_VECTOR_FORMATS = {"wmf", "emf", "x-wmf", "x-emf"}


def _extract_pptx_images(pptx_path: Path, output_dir: Path) -> dict[int, list[Path]]:
    """Extract images from PPTX using python-pptx with group recursion."""
    from pptx import Presentation as PptxPresentation  # noqa: PLC0415
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

    prs = PptxPresentation(str(pptx_path))
    fig_dir = output_dir / "figures"
    fig_dir.mkdir(parents=True, exist_ok=True)

    images: dict[int, list[Path]] = {}
    fig_count = 0
    wmf_skipped = 0
    seen_hashes: dict[str, Path] = {}  # hash -> filepath (for dedup)
    dedup_count = 0

    def process_shape(shape: object, slide_num: int) -> None:
        nonlocal fig_count, wmf_skipped, dedup_count
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # type: ignore[attr-defined]
            for grouped in shape.shapes:  # type: ignore[attr-defined]
                process_shape(grouped, slide_num)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # type: ignore[attr-defined]
            try:
                image = shape.image  # type: ignore[attr-defined]
                content_type = image.content_type
                if any(vf in content_type for vf in _VECTOR_FORMATS):
                    wmf_skipped += 1
                    return
                blob = image.blob
                blob_hash = hashlib.md5(blob).hexdigest()
                if blob_hash in seen_hashes:
                    images.setdefault(slide_num, []).append(seen_hashes[blob_hash])
                    dedup_count += 1
                    return
                ext = content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                elif ext not in ("png", "jpg", "gif", "bmp", "tiff", "webp"):
                    ext = "png"
                filename = f"figure_{fig_count}.{ext}"
                filepath = fig_dir / filename
                filepath.write_bytes(blob)
                seen_hashes[blob_hash] = filepath
                images.setdefault(slide_num, []).append(filepath)
                fig_count += 1
            except Exception:
                logger.warning("Failed to extract image from slide %d", slide_num)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            process_shape(shape, slide_idx + 1)

    total = fig_count + dedup_count
    logger.info("Extracted %d unique image(s) from %d total (%d duplicates removed)", fig_count, total, dedup_count)
    if wmf_skipped > 0:
        logger.warning("SKIPPED %d WMF/EMF image(s) (vector format not supported by PIL)", wmf_skipped)
    return images


def _map_images_to_items(
    doc: object, slide_images: dict[int, list[Path]]
) -> tuple[dict[str, str], list[Path], list[str]]:
    """Map python-pptx extracted images to Docling PictureItems by slide + order."""
    from docling_core.types.doc.document import PictureItem  # noqa: PLC0415

    figure_map: dict[str, str] = {}
    image_paths: list[Path] = []
    item_refs: list[str] = []
    picture_idx_per_slide: dict[int, int] = {}

    for item, _ in doc.iterate_items():  # type: ignore[attr-defined]
        if isinstance(item, PictureItem):
            prov = getattr(item, "prov", None)
            slide_num = prov[0].page_no if prov else 0
            idx = picture_idx_per_slide.get(slide_num, 0)
            picture_idx_per_slide[slide_num] = idx + 1

            slide_imgs = slide_images.get(slide_num, [])
            if idx < len(slide_imgs):
                figure_map[item.self_ref] = f"figures/{slide_imgs[idx].name}"
                image_paths.append(slide_imgs[idx])
                item_refs.append(item.self_ref)

    return figure_map, image_paths, item_refs


class PptxConverter(BaseConverter):
    """PPTX conversion: Docling text + python-pptx images + VLM descriptions."""

    def convert(self) -> None:
        from docling.datamodel.base_models import InputFormat  # noqa: PLC0415
        from docling.document_converter import DocumentConverter, PowerpointFormatOption  # noqa: PLC0415

        from tracing import trace_span  # noqa: PLC0415

        self.ensure_output_dir()

        converter = DocumentConverter(
            allowed_formats=[InputFormat.PPTX],
            format_options={InputFormat.PPTX: PowerpointFormatOption()},
        )

        with trace_span("docling.convert_pptx", file=self.source.name):
            logger.info("Converting %s (PPTX native)", self.source.name)
            result = converter.convert(str(self.source))
            logger.info("Status: %s", result.status)

        doc = result.document
        contexts = collect_floating_contexts(doc)
        context_blocks = _build_context_blocks(contexts)
        artifacts = FloatingArtifacts()

        fig_count = 0
        if self.options.figures:
            slide_images = _extract_pptx_images(self.source, self.output_dir)
            figure_map, image_paths, item_refs = _map_images_to_items(doc, slide_images)
            artifacts.figure_paths = figure_map
            fig_count = len(image_paths)
            artifacts.figure_descriptions = self.describe_figures(
                image_paths, item_refs, "vlm.describe_pptx_images", context_blocks
            )

        title = extract_title(doc)
        page_md = build_document_markdown(doc, artifacts, contexts, title=title, paginated=True)
        self.write_document_md(page_md)

        if fig_count > 0:
            catalog = build_images_catalog(doc, artifacts, contexts)
            (self.output_dir / "images.md").write_text(catalog)

        if self.options.all_formats:
            self.write_all_formats(doc)

        if title:
            logger.info("Title: %s", title)
        self.print_summary(
            fig_count=fig_count,
            captions_used=self.options.captions_enabled and fig_count > 0,
            desc_count=len(artifacts.figure_descriptions),
        )


def _build_context_blocks(contexts: dict[str, FloatingContext]) -> dict[str, str]:
    from doc_convert.providers import build_context_block  # noqa: PLC0415

    return {ref: build_context_block(caption=ctx.caption, mention=ctx.mention) for ref, ctx in contexts.items()}
