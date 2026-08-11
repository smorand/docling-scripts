"""PPTX whole-slide screenshot rendering, notes extraction, and slide VLM analysis.

Slides are rendered faithfully by converting the PPTX to PDF via a headless
LibreOffice (``soffice``) process, then rasterizing each PDF page with
``pypdfium2`` (already a project dependency, used elsewhere for PDF metadata).
This avoids re-implementing PowerPoint layout/rendering by hand, which would
never match fonts, gradients, SmartArt, or animations-at-rest reliably.

Each rendered slide image is then sent, on its own, to a capable multimodal
LLM (default: IBM ICA ``claude-sonnet-4-6``) for a detailed visual
interpretation. This is intentionally independent from the extracted text +
figures pipeline (see ``converters/pptx.py``): the two are combined only at
markdown-assembly time, as separate labelled sections, so the LLM consuming
``document.md`` can tell "what was mechanically extracted" apart from "what a
vision model saw on the rendered slide".
"""

from __future__ import annotations

import logging
import re
import shutil
import subprocess
import tempfile
import uuid
from dataclasses import dataclass
from pathlib import Path

import typer

from config import Settings  # noqa: TC001
from doc_convert.providers import (
    DEFAULT_LLM_CONCURRENCY,
    DEFAULT_PPTX_SLIDE_VLM,
    parse_external_llm,
    require_api_key,
)
from logging_config import console

logger = logging.getLogger(__name__)

# Matches the PDF pipeline's images_scale convention (sharper crops for VLM input).
_SLIDE_IMAGE_SCALE = 2.0

# soffice conversion can be slow on large decks with embedded media.
_SOFFICE_TIMEOUT_SECONDS = 300

# HTTP status meaning "model slug not found/served" on the OpenAI-compatible endpoint.
_HTTP_NOT_FOUND = 404

# Threads are safe on this path because it touches no docling pipeline, no torch
# and no global state: httpx.Client is documented as thread-safe, and image_prep
# writes its temporaries through tempfile.mkstemp. The default worker count
# lives in providers.py next to the other provider-facing defaults.

# Concurrency raises the odds of hitting a rate limit, and before retries a
# single 429 silently dropped that slide's analysis from document.md. Retry
# transient statuses so parallelism cannot degrade output quality.
_MAX_SLIDE_ATTEMPTS = 4

_SLIDE_ANALYSIS_PROMPT = (
    "You are analyzing a screenshot of a single PowerPoint slide, rendered exactly "
    "as it would display. Write a detailed markdown interpretation of what this "
    "slide visually communicates.\n\n"
    "Cover, in order:\n"
    "1. **Purpose**: what the slide is trying to communicate.\n"
    "2. **Layout**: how text and visuals are arranged, and the visual hierarchy "
    "(what stands out first, second, etc.).\n"
    "3. **Visual elements**: describe every chart, diagram, schema, screenshot, "
    "photo, icon, table, callout, arrow, or annotation. For charts/diagrams, "
    "explain what they show and how the parts relate, not just that they exist.\n"
    "4. **Key text**: mention critical text (titles, key numbers, labels) that "
    "matters for interpreting the slide, without transcribing everything verbatim.\n\n"
    "Rules:\n"
    "- Focus on visual meaning and interpretation, not a literal transcription.\n"
    "- Never say the image is unreadable or that you cannot see it, unless it "
    "genuinely is (blank/corrupted).\n"
    "- This response will be embedded as a subsection inside a larger document. "
    "Do NOT start with a top-level title for the whole response (no '#' or '##' "
    "heading naming/labelling the analysis itself, e.g. do not write 'Slide "
    "Analysis' or repeat the slide title as a heading). Start directly with the "
    "content (plain text or a bold lead-in). If you use sub-headings for the 4 "
    "points above, use '####' or deeper.\n"
    "- Output markdown only. No preamble, no closing remarks."
)

_SLIDE_USER_PROMPT = "Analyze this PowerPoint slide screenshot in detail."

_HEADING_RE = re.compile(r"^(#{1,6})(\s+\S)")


def _normalize_heading_levels(text: str, min_level: int = 4) -> str:
    """Demote any markdown heading in ``text`` below ``min_level``.

    The VLM is instructed not to emit a top-level heading, but models don't
    always comply (e.g. opening with ``## Slide Analysis``). Since this text
    is spliced in as content under our own ``### Visual Interpretation``
    subsection, an errant ``##``/``###`` would visually collide with our
    ``## Slide N`` / ``### Visual Interpretation`` structure. This is a
    deterministic safety net independent of prompt compliance.
    """
    lines = text.splitlines()
    for i, line in enumerate(lines):
        m = _HEADING_RE.match(line)
        if m and len(m.group(1)) < min_level:
            lines[i] = "#" * min_level + line[len(m.group(1)) :]
    return "\n".join(lines)


@dataclass(frozen=True)
class _SlideAttempt:
    """Outcome of one slide analysis call."""

    description: str = ""
    retryable: bool = False
    reason: str = ""


@dataclass(frozen=True)
class SlideVisualAnalysis:
    """Result of sending one rendered slide screenshot to the VLM."""

    slide_number: int
    image_path: str
    description: str


@dataclass(frozen=True)
class SlideRenderArtifacts:
    """Everything gathered for the 'visual interpretation' pass, per slide."""

    image_paths_by_slide: dict[int, str]
    analyses_by_slide: dict[int, SlideVisualAnalysis]
    notes_by_slide: dict[int, str]
    hidden_slides: set[int]


def extract_slide_notes(pptx_path: Path) -> dict[int, str]:
    """Extract speaker notes text per slide (1-indexed), skipping empty notes."""
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(pptx_path))
    notes: dict[int, str] = {}
    for idx, slide in enumerate(prs.slides, 1):
        if not slide.has_notes_slide:
            continue
        text = (slide.notes_slide.notes_text_frame.text or "").strip()
        if text:
            notes[idx] = text
    return notes


def extract_hidden_slide_numbers(pptx_path: Path) -> set[int]:
    """Return the 1-indexed slide numbers marked as hidden in PowerPoint.

    Hidden slides still get a full extracted-content + screenshot + notes
    treatment (nothing is skipped), but callers can annotate them so a
    downstream reader knows the slide would not appear in an actual
    presentation.
    """
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(pptx_path))
    return {idx for idx, slide in enumerate(prs.slides, 1) if slide._element.get("show") == "0"}


def render_pptx_slides(pptx_path: Path, output_dir: Path) -> dict[int, str]:
    """Render every slide to a PNG screenshot under ``<output_dir>/slides/``.

    Returns {slide_number: relative_path}. Slide numbering is 1-indexed and
    matches python-pptx / docling page numbering for this same file.
    """
    import pypdfium2  # noqa: PLC0415

    slides_dir = output_dir / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="doc_convert_pptx_render_") as tmp:
        pdf_path = _pptx_to_pdf(pptx_path, Path(tmp))
        pdf = pypdfium2.PdfDocument(str(pdf_path))
        try:
            rendered: dict[int, str] = {}
            for idx, page in enumerate(pdf, 1):
                bitmap = page.render(scale=_SLIDE_IMAGE_SCALE)
                image = bitmap.to_pil().convert("RGB")
                out_path = slides_dir / f"slide_{idx:03d}.png"
                image.save(out_path)
                rendered[idx] = f"slides/{out_path.name}"
        finally:
            pdf.close()
    logger.info("Rendered %d slide screenshot(s) to %s", len(rendered), slides_dir)
    return rendered


def analyze_slide_images(
    image_paths_by_slide: dict[int, str],
    output_dir: Path,
    settings: Settings,
    *,
    llm: str = DEFAULT_PPTX_SLIDE_VLM,
    concurrency: int = DEFAULT_LLM_CONCURRENCY,
) -> dict[int, SlideVisualAnalysis]:
    """Send each rendered slide screenshot to the VLM for a visual interpretation.

    Slides are analysed ``concurrency`` at a time via the shared vision primitive.
    Results are keyed by slide number, so completion order never affects the
    assembled document. A slide that fails every retry is simply absent.
    """
    from doc_convert.vision_llm import describe_image, make_client, map_concurrent  # noqa: PLC0415

    provider, model = parse_external_llm(llm)
    api_key = require_api_key(provider, settings)

    pending: list[tuple[int, str]] = []
    for slide_number, rel_path in sorted(image_paths_by_slide.items()):
        if (output_dir / rel_path).exists():
            pending.append((slide_number, rel_path))
        else:
            logger.warning("Slide screenshot not found, skipping: %s", output_dir / rel_path)

    if not pending:
        return {}

    total = len(pending)
    workers = max(1, min(concurrency, total))
    logger.info("Analyzing %d slide screenshot(s) with %s/%s (%d at a time)", total, provider, model, workers)

    def report(done: int, item: tuple[int, str], description: str, elapsed: float) -> None:
        if description:
            logger.info("Analyzed slide %d in %.1fs (%d/%d done)", item[0], elapsed, done, total)
        else:
            logger.warning("Skipping slide %d visual analysis (API error)", item[0])

    with make_client(settings, workers) as client:

        def work(item: tuple[int, str]) -> str:
            _slide_number, rel_path = item
            return describe_image(
                output_dir / rel_path,
                _SLIDE_USER_PROMPT,
                provider,
                model,
                settings,
                api_key,
                client,
                system=_SLIDE_ANALYSIS_PROMPT,
            )

        descriptions = map_concurrent(pending, work, workers, what="slide", on_done=report)

    return {
        slide_number: SlideVisualAnalysis(
            slide_number=slide_number,
            image_path=rel_path,
            description=_normalize_heading_levels(description.strip()),
        )
        for (slide_number, rel_path), description in zip(pending, descriptions, strict=True)
        if description
    }


def build_slide_render_artifacts(
    pptx_path: Path,
    output_dir: Path,
    settings: Settings,
    *,
    llm: str = DEFAULT_PPTX_SLIDE_VLM,
    concurrency: int = DEFAULT_LLM_CONCURRENCY,
) -> SlideRenderArtifacts:
    """Render slides, run the VLM visual analysis, and extract speaker notes."""
    # Resolve to absolute path immediately so that relative paths stay valid
    # regardless of any cwd changes during the conversion pipeline.
    output_dir = output_dir.resolve()
    image_paths_by_slide = render_pptx_slides(pptx_path, output_dir)
    analyses_by_slide = analyze_slide_images(
        image_paths_by_slide, output_dir, settings, llm=llm, concurrency=concurrency
    )
    notes_by_slide = extract_slide_notes(pptx_path)
    hidden_slides = extract_hidden_slide_numbers(pptx_path)
    return SlideRenderArtifacts(
        image_paths_by_slide=image_paths_by_slide,
        analyses_by_slide=analyses_by_slide,
        notes_by_slide=notes_by_slide,
        hidden_slides=hidden_slides,
    )


def _pptx_to_pdf(pptx_path: Path, workdir: Path) -> Path:
    """Convert a PPTX to PDF via headless LibreOffice, for faithful page rasterization.

    Uses a unique ``UserInstallation`` profile dir so concurrent conversions
    (e.g. under ``-P`` batch mode) never contend on the same LibreOffice lock.
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        console.print(
            "[red]soffice (LibreOffice) is required to render PPTX slide screenshots. "
            "Install it with: brew install libreoffice[/red]"
        )
        raise typer.Exit(1)

    profile_dir = workdir / f"lo_profile_{uuid.uuid4().hex}"
    # ExportHiddenSlides=true is required: LibreOffice's PDF export silently
    # drops slides marked hidden in PowerPoint by default, which would shift
    # every subsequent PDF page number out of sync with the docling/python-pptx
    # slide numbering (used for text content and speaker notes). Keeping every
    # slide in the PDF, in original order, is what keeps the 3 sections of a
    # given "## Slide N" aligned to the same physical slide.
    convert_filter = 'pdf:impress_pdf_Export:{"ExportHiddenSlides":{"type":"boolean","value":"true"}}'
    cmd = [
        soffice,
        "--headless",
        "--norestore",
        f"-env:UserInstallation=file://{profile_dir}",
        "--convert-to",
        convert_filter,
        "--outdir",
        str(workdir),
        str(pptx_path),
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=_SOFFICE_TIMEOUT_SECONDS, check=False)
    except subprocess.TimeoutExpired as exc:
        console.print(f"[red]soffice timed out converting {pptx_path.name} to PDF[/red]")
        raise typer.Exit(1) from exc

    if result.returncode != 0:
        console.print(f"[red]soffice failed to convert {pptx_path.name} to PDF: {result.stderr.strip()}[/red]")
        raise typer.Exit(1)

    pdf_path = workdir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        console.print(f"[red]soffice did not produce the expected PDF at {pdf_path}[/red]")
        raise typer.Exit(1)
    return pdf_path
