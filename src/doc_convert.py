"""Unified document converter using Docling.

Converts PDF, images, DOCX, XLSX, and Google Docs/Sheets to markdown.
PDF uses local models by default (full extraction with figures, tables, OCR).
Use --gemini for quick conversion via Gemini API.
"""

from __future__ import annotations

import json
import logging
import os
import re
import tempfile
from enum import Enum
from pathlib import Path

import httpx
import pypdfium2
import typer
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import (
    PdfPipelineOptions,
    PictureDescriptionVlmEngineOptions,
    VlmPipelineOptions,
)
from docling.datamodel.pipeline_options_vlm_model import (
    ApiVlmOptions,
    ResponseFormat,
)
from docling.document_converter import (
    DocumentConverter,
    ExcelFormatOption,
    FormatOption,
    ImageFormatOption,
    PdfFormatOption,
    PowerpointFormatOption,
    WordFormatOption,
)
from docling.pipeline.vlm_pipeline import VlmPipeline
from docling_core.types.doc import ImageRefMode
from docling_core.types.doc.document import PictureItem, TableItem
from docling_core.types.doc.labels import DocItemLabel

from config import Settings
from logging_config import console, setup_logging
from tracing import configure_tracing, trace_span

logger = logging.getLogger(__name__)

app = typer.Typer(
    help="Convert documents to markdown using Docling.",
    no_args_is_help=True,
)


# ── Enums ─────────────────────────────────────────────────────────────────────


class VlmPreset(str, Enum):
    smolvlm = "smolvlm"
    granite_vision = "granite_vision"
    pixtral = "pixtral"
    qwen = "qwen"


PRESET_REPO_IDS: dict[str, str] = {
    "smolvlm": "HuggingFaceTB/SmolVLM-256M-Instruct",
    "granite_vision": "ibm-granite/granite-vision-3.3-2b",
    "pixtral": "mistral-community/pixtral-12b",
    "qwen": "Qwen/Qwen2.5-VL-3B-Instruct",
}


# ── Format detection ──────────────────────────────────────────────────────────

EXT_TO_FORMAT: dict[str, InputFormat] = {
    ".pdf": InputFormat.PDF,
    ".docx": InputFormat.DOCX,
    ".xlsx": InputFormat.XLSX,
    ".pptx": InputFormat.PPTX,
    ".pptm": InputFormat.PPTX,
    ".potx": InputFormat.PPTX,
    ".ppsx": InputFormat.PPTX,
    ".jpg": InputFormat.IMAGE,
    ".jpeg": InputFormat.IMAGE,
    ".png": InputFormat.IMAGE,
    ".tiff": InputFormat.IMAGE,
    ".tif": InputFormat.IMAGE,
    ".bmp": InputFormat.IMAGE,
    ".webp": InputFormat.IMAGE,
}


def detect_format(path: Path) -> InputFormat:
    """Detect input format from file extension."""
    ext = path.suffix.lower()
    fmt = EXT_TO_FORMAT.get(ext)
    if fmt is None:
        console.print(f"[red]Unsupported extension '{ext}'[/red]")
        console.print(f"Supported: {', '.join(sorted(EXT_TO_FORMAT.keys()))}")
        raise typer.Exit(1)
    return fmt


# ── Google Docs/Sheets support ────────────────────────────────────────────────

GOOGLE_DOC_RE = re.compile(r"https://docs\.google\.com/document/d/([a-zA-Z0-9_-]+)")
GOOGLE_SHEET_RE = re.compile(r"https://docs\.google\.com/spreadsheets/d/([a-zA-Z0-9_-]+)")
MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
DRIVE_EXPORT_URL = "https://www.googleapis.com/drive/v3/files/{file_id}/export"
DRIVE_FILE_URL = "https://www.googleapis.com/drive/v3/files/{file_id}"
DRIVE_SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]


def is_google_url(source: str) -> bool:
    """Check if the input is a Google Docs or Sheets URL."""
    return bool(GOOGLE_DOC_RE.search(source) or GOOGLE_SHEET_RE.search(source))


def _load_google_credentials(settings: Settings) -> str:
    """Load Google credentials and return an access token."""
    creds_path = settings.google_credentials
    if not creds_path:
        console.print("[red]GOOGLE_CREDENTIALS env var is required for Google Docs/Sheets[/red]")
        raise typer.Exit(1)

    creds_file = Path(os.path.expandvars(creds_path)).expanduser()
    if not creds_file.exists():
        console.print(f"[red]Credentials file not found: {creds_file}[/red]")
        raise typer.Exit(1)

    from google.oauth2 import credentials as user_credentials  # noqa: PLC0415
    from google.oauth2 import service_account  # noqa: PLC0415

    creds_data = json.loads(creds_file.read_text())

    cred_type = creds_data.get("type", "")
    if cred_type == "service_account":
        creds = service_account.Credentials.from_service_account_file(str(creds_file), scopes=DRIVE_SCOPES)
    elif cred_type == "authorized_user":
        creds = user_credentials.Credentials.from_authorized_user_file(str(creds_file), scopes=DRIVE_SCOPES)
    else:
        console.print(f"[red]Unsupported credential type '{cred_type}'[/red]")
        raise typer.Exit(1)

    from google.auth.transport.requests import Request as AuthRequest  # noqa: PLC0415

    if not creds.valid:
        creds.refresh(AuthRequest())
    return creds.token


def download_google_doc(url: str, settings: Settings) -> tuple[Path, str, InputFormat]:
    """Download a Google Doc/Sheet to a temp file."""
    doc_match = GOOGLE_DOC_RE.search(url)
    sheet_match = GOOGLE_SHEET_RE.search(url)

    if doc_match:
        file_id = doc_match.group(1)
        mime_type, suffix = MIME_DOCX, ".docx"
        fmt, kind = InputFormat.DOCX, "Google Doc"
    elif sheet_match:
        file_id = sheet_match.group(1)
        mime_type, suffix = MIME_XLSX, ".xlsx"
        fmt, kind = InputFormat.XLSX, "Google Sheet"
    else:
        console.print(f"[red]Not a recognized Google Docs/Sheets URL: {url}[/red]")
        raise typer.Exit(1)

    token = _load_google_credentials(settings)
    headers = {"Authorization": f"Bearer {token}"}

    with trace_span("google.download", kind=kind, file_id=file_id), httpx.Client(timeout=60.0) as client:
        title_resp = client.get(
            DRIVE_FILE_URL.format(file_id=file_id),
            headers=headers,
            params={"fields": "name"},
        )
        title = title_resp.json().get("name", file_id) if title_resp.is_success else file_id
        logger.info("Downloading %s: %s", kind, title)

        resp = client.get(
            DRIVE_EXPORT_URL.format(file_id=file_id),
            headers=headers,
            params={"mimeType": mime_type},
        )
        if not resp.is_success:
            console.print(f"[red]Failed to export {kind} (HTTP {resp.status_code})[/red]")
            raise typer.Exit(1)

    tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)  # noqa: SIM115
    tmp.write(resp.content)
    tmp.close()
    return Path(tmp.name), title, fmt


# ── External LLM providers ───────────────────────────────────────────────────

EXTERNAL_LLM_PROMPT = (
    "Convert this document page to well-structured markdown. "
    "Extract ALL text precisely.\n\n"
    "For administrative documents, clearly identify and highlight:\n"
    "- Personal identifiers (passport numbers, ID numbers, client numbers, "
    "social security numbers)\n"
    "- Credentials (login, passwords, access codes)\n"
    "- Dates (issue dates, expiry dates, deadlines, birth dates)\n"
    "- Locations (addresses, cities, countries)\n"
    "- People and their roles (signatories, mandated persons, "
    "representatives, beneficiaries)\n"
    "- Financial amounts (costs, revenues, taxes, fees, totals "
    "with currency)\n"
    "- Reference numbers (invoice numbers, contract numbers, case numbers)\n\n"
    "Format these as bold or in a clearly labeled section. "
    "Do not miss any text. Output only the bare markdown."
)

PROVIDER_URLS: dict[str, str] = {
    "google": "https://generativelanguage.googleapis.com/v1beta/openai/chat/completions",
    "openrouter": "https://openrouter.ai/api/v1/chat/completions",
}


def parse_external_llm(value: str) -> tuple[str, str]:
    """Parse 'provider/model' into (provider, model).

    Examples:
        google/gemini-3-flash-preview -> ("google", "gemini-3-flash-preview")
        openrouter/google/gemini-3-pro-preview -> ("openrouter", "google/gemini-3-pro-preview")
    """
    for provider in PROVIDER_URLS:
        prefix = f"{provider}/"
        if value.startswith(prefix):
            model = value[len(prefix) :]
            if not model:
                console.print(f"[red]Missing model name after '{prefix}'[/red]")
                raise typer.Exit(1)
            return provider, model
    supported = ", ".join(f"{p}/<model>" for p in PROVIDER_URLS)
    console.print(f"[red]Unknown provider in '{value}'. Supported: {supported}[/red]")
    raise typer.Exit(1)


def _require_api_key(provider: str, settings: Settings) -> str:
    """Get the API key for a provider, or exit with error."""
    if provider == "google":
        if not settings.google_api_key:
            console.print("[red]GOOGLE_API_KEY env var is required for google/ provider[/red]")
            raise typer.Exit(1)
        return settings.google_api_key
    if provider == "openrouter":
        if not settings.openrouter_api_key:
            console.print("[red]OPENROUTER_API_KEY env var is required for openrouter/ provider[/red]")
            raise typer.Exit(1)
        return settings.openrouter_api_key
    console.print(f"[red]Unknown provider: {provider}[/red]")
    raise typer.Exit(1)


def build_external_vlm_options(provider: str, model: str, settings: Settings) -> ApiVlmOptions:
    """Build VLM options for an external LLM provider."""
    api_key = _require_api_key(provider, settings)
    return ApiVlmOptions(
        url=PROVIDER_URLS[provider],
        params={"model": model, "max_tokens": settings.llm_max_tokens},
        headers={"Authorization": f"Bearer {api_key}"},
        prompt=EXTERNAL_LLM_PROMPT,
        scale=2.0,
        timeout=settings.llm_timeout,
        response_format=ResponseFormat.MARKDOWN,
        temperature=0.0,
    )


# ── PDF local extraction helpers ──────────────────────────────────────────────


def get_document_title(doc: object, pdf_path: str) -> str:
    """Extract document title from docling document or PDF metadata."""
    for t in doc.texts:  # type: ignore[attr-defined]
        if t.label == DocItemLabel.TITLE:
            return t.text  # type: ignore[no-any-return]
    try:
        pdf = pypdfium2.PdfDocument(pdf_path)
        meta = pdf.get_metadata_dict()
        title = meta.get("Title", "")
        if title:
            return title  # type: ignore[no-any-return]
    except Exception:
        pass
    return ""


def get_pdf_metadata(pdf_path: str) -> dict[str, str]:
    """Extract metadata from PDF file."""
    meta: dict[str, str] = {}
    try:
        pdf = pypdfium2.PdfDocument(pdf_path)
        raw = pdf.get_metadata_dict()
        for key in ("Author", "Subject", "Keywords", "Creator", "CreationDate"):
            val = raw.get(key, "")
            if val:
                meta[key] = val
        meta["Pages"] = str(len(pdf))
        meta["File"] = Path(pdf_path).name
    except Exception:
        meta["File"] = Path(pdf_path).name
    return meta


def get_vlm_description(item: PictureItem) -> str:
    """Extract VLM-generated description from item metadata."""
    meta = getattr(item, "meta", None)
    if meta is None:
        return ""
    desc = getattr(meta, "description", None)
    if desc is None:
        return ""
    return getattr(desc, "text", "") or ""


def get_picture_classification(item: PictureItem) -> str:
    """Extract picture classification label."""
    meta = getattr(item, "meta", None)
    if meta is None:
        return ""
    cls_info = getattr(meta, "classification", None)
    if cls_info is None:
        return ""
    return getattr(cls_info, "label", "") or ""


def build_page_annotated_markdown(  # noqa: PLR0912
    doc: object,
    figure_map: dict[str, str],
    title: str = "",
    pdf_meta: dict[str, str] | None = None,
    figure_descriptions: dict[str, str] | None = None,
) -> str:
    """Build markdown with page number annotations."""
    lines: list[str] = []
    current_page = None

    if title:
        lines.append(f"# {title}\n")

    if pdf_meta:
        lines.append("| | |")
        lines.append("|---|---|")
        for key, val in pdf_meta.items():
            lines.append(f"| **{key}** | {val} |")
        lines.append("")

    for item, _ in doc.iterate_items():  # type: ignore[attr-defined]
        prov = getattr(item, "prov", None)
        if prov:
            page_no = prov[0].page_no
            if page_no != current_page:
                if current_page is not None:
                    lines.append("")
                lines.append(f"---\n*[Page {page_no}]*\n")
                current_page = page_no

        if isinstance(item, TableItem):
            df = item.export_to_dataframe(doc)
            lines.append(df.to_markdown(index=False))
            caption = item.caption_text(doc)
            if caption:
                lines.append(f"\n*{caption}*")
            lines.append("")
        elif isinstance(item, PictureItem):
            caption = item.caption_text(doc)
            description = (figure_descriptions or {}).get(item.self_ref, "") or get_vlm_description(item)
            fig_path = figure_map.get(item.self_ref, "")

            lines.append(f"[Figure: {caption}]" if caption else "[Figure]")
            if fig_path:
                lines.append(f"![figure]({fig_path})")
            if description:
                lines.append(f"\n> {description}")
            lines.append("")
        else:
            text = getattr(item, "text", None)
            if text:
                label = getattr(item, "label", "")
                if "section_header" in str(label).lower():
                    level = getattr(item, "level", 1)
                    prefix = "#" * min(level + 1, 6)
                    lines.append(f"{prefix} {text}\n")
                elif "list_item" in str(label).lower():
                    lines.append(f"- {text}")
                else:
                    lines.append(f"{text}\n")

    return "\n".join(lines)


def build_images_catalog(doc: object, figure_map: dict[str, str]) -> str:
    """Build an images.md catalog (similar to pdf-extractor output)."""
    lines: list[str] = ["# Image Catalog\n"]

    for item, _ in doc.iterate_items():  # type: ignore[attr-defined]
        if not isinstance(item, PictureItem):
            continue
        fig_path = figure_map.get(item.self_ref, "")
        if not fig_path:
            continue

        description = get_vlm_description(item)
        caption = item.caption_text(doc)
        classification = get_picture_classification(item)

        lines.append(f"## {Path(fig_path).name}\n")
        lines.append(f"- **Path:** {fig_path}")
        if classification:
            lines.append(f"- **Type:** {classification}")
        if caption:
            lines.append(f"- **Caption:** {caption}")
        if description:
            lines.append(f"- **Description:** {description}")
        lines.append(f"\n![{Path(fig_path).stem}]({fig_path})\n")

    return "\n".join(lines)


# ── PPTX image extraction ────────────────────────────────────────────────────


def extract_pptx_images(pptx_path: Path, output_dir: Path) -> dict[int, list[Path]]:
    """Extract images from PPTX using python-pptx with group recursion.

    Returns a dict mapping slide number (1-based) to list of extracted image paths.
    """
    from pptx import Presentation as PptxPresentation  # noqa: PLC0415
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

    prs = PptxPresentation(str(pptx_path))
    fig_dir = output_dir / "figures"
    fig_dir.mkdir(parents=True, exist_ok=True)

    images: dict[int, list[Path]] = {}
    fig_count = 0

    def process_shape(shape: object, slide_num: int) -> None:
        nonlocal fig_count
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # type: ignore[attr-defined]
            for grouped in shape.shapes:  # type: ignore[attr-defined]
                process_shape(grouped, slide_num)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # type: ignore[attr-defined]
            try:
                image = shape.image  # type: ignore[attr-defined]
                content_type = image.content_type
                ext = content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                elif ext not in ("png", "jpg", "gif", "bmp", "tiff", "webp"):
                    ext = "png"
                filename = f"slide{slide_num}_fig{fig_count}.{ext}"
                filepath = fig_dir / filename
                filepath.write_bytes(image.blob)
                images.setdefault(slide_num, []).append(filepath)
                fig_count += 1
            except Exception:
                logger.warning("Failed to extract image from slide %d", slide_num)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            process_shape(shape, slide_idx + 1)

    logger.info("Extracted %d image(s) from %d slide(s)", fig_count, len(prs.slides))
    return images


def describe_images_with_vlm(image_paths: list[Path], vlm_preset: str, models_path: Path) -> list[str]:
    """Describe images using local VLM engine (offline, uses models_path)."""
    from docling.datamodel.accelerator_options import AcceleratorOptions  # noqa: PLC0415
    from docling.models.stages.picture_description.picture_description_vlm_engine_model import (  # noqa: PLC0415
        PictureDescriptionVlmEngineModel,
    )
    from PIL import Image  # noqa: PLC0415

    if not models_path.exists():
        console.print(f"[red]Models directory not found: {models_path}[/red]")
        console.print("Run [bold]doc-convert download-models[/bold] first.")
        raise typer.Exit(1)

    options = PictureDescriptionVlmEngineOptions.from_preset(vlm_preset)
    model = PictureDescriptionVlmEngineModel(
        enabled=True,
        enable_remote_services=False,
        artifacts_path=str(models_path),
        options=options,
        accelerator_options=AcceleratorOptions(),
    )

    logger.info("Describing %d image(s) with local VLM (%s)", len(image_paths), vlm_preset)
    pil_images = [Image.open(p) for p in image_paths]
    raw = list(model._annotate_images(pil_images))
    return [re.sub(r"<end_of_utteranc\w*>?", "", d).strip() for d in raw]


def describe_images_with_external_llm(
    image_paths: list[Path], provider: str, model: str, settings: Settings
) -> list[str]:
    """Describe images using an external LLM provider."""
    logger.info("Describing %d image(s) with %s/%s", len(image_paths), provider, model)
    descriptions: list[str] = []
    for img_path in image_paths:
        md = convert_with_external_llm(img_path, InputFormat.IMAGE, provider, model, settings)
        descriptions.append(md.strip())
    return descriptions


# ── Conversion functions ──────────────────────────────────────────────────────


def convert_pdf_local(
    pdf_path: Path,
    output_dir: Path,
    *,
    do_ocr: bool,
    vlm: bool,
    vlm_preset: str,
    figures: bool,
    all_formats: bool,
    models_path: Path,
) -> None:
    """Full PDF extraction with local models."""
    output_dir.mkdir(parents=True, exist_ok=True)

    pipeline_options = PdfPipelineOptions(
        do_ocr=do_ocr,
        do_table_structure=True,
        generate_page_images=figures,
        generate_picture_images=figures,
        do_picture_description=vlm and figures,
        do_picture_classification=figures,
        artifacts_path=str(models_path),
    )

    if vlm:
        pipeline_options.picture_description_options = PictureDescriptionVlmEngineOptions.from_preset(vlm_preset)
        logger.info("VLM picture description enabled (preset: %s)", vlm_preset)

    converter = DocumentConverter(format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)})

    with trace_span("docling.convert_pdf_local", file=pdf_path.name):
        logger.info("Converting %s (local pipeline)", pdf_path.name)
        result = converter.convert(str(pdf_path))
        logger.info("Status: %s", result.status)

    doc = result.document

    # Extract figures (skip entirely when --no-figures)
    figure_map: dict[str, str] = {}
    fig_count = 0
    if figures:
        fig_dir = output_dir / "figures"
        fig_dir.mkdir(exist_ok=True)
        for item, _ in doc.iterate_items():
            if isinstance(item, PictureItem):
                img = item.get_image(doc)
                if img:
                    filename = f"figure_{fig_count}.png"
                    img.save(fig_dir / filename)
                    figure_map[item.self_ref] = f"figures/{filename}"
                    fig_count += 1

    title = get_document_title(doc, str(pdf_path))
    pdf_meta = get_pdf_metadata(str(pdf_path))

    page_md = build_page_annotated_markdown(doc, figure_map, title, pdf_meta)
    (output_dir / "document.md").write_text(page_md)

    if fig_count > 0:
        catalog = build_images_catalog(doc, figure_map)
        (output_dir / "images.md").write_text(catalog)

    if all_formats:
        (output_dir / "output.md").write_text(doc.export_to_markdown())
        (output_dir / "output_embedded.md").write_text(doc.export_to_markdown(image_mode=ImageRefMode.EMBEDDED))
        (output_dir / "output.json").write_text(json.dumps(doc.export_to_dict(), indent=2, default=str))
        (output_dir / "output.txt").write_text(doc.export_to_text())
        (output_dir / "output.html").write_text(doc.export_to_html())

    if title:
        logger.info("Title: %s", title)
    console.print(f"[green]Output:[/green] {output_dir}/")
    console.print("  document.md  (page-annotated markdown)")
    if fig_count > 0:
        console.print("  images.md    (image catalog)")
        console.print(f"  figures/     ({fig_count} figure(s))")
    if all_formats:
        console.print("  output.*     (md, html, json, txt, embedded)")


def convert_with_external_llm(doc_path: Path, fmt: InputFormat, provider: str, model: str, settings: Settings) -> str:
    """Convert PDF or image via external LLM VLM pipeline."""
    vlm_opts = build_external_vlm_options(provider, model, settings)
    pipeline_options = VlmPipelineOptions(
        enable_remote_services=True,
        vlm_options=vlm_opts,
    )

    format_options: dict[InputFormat, FormatOption] = {}
    if fmt == InputFormat.PDF:
        format_options[InputFormat.PDF] = PdfFormatOption(
            pipeline_cls=VlmPipeline,
            pipeline_options=pipeline_options,
        )
    else:
        format_options[InputFormat.IMAGE] = ImageFormatOption(
            pipeline_cls=VlmPipeline,
            pipeline_options=pipeline_options,
        )

    converter = DocumentConverter(
        allowed_formats=[fmt],
        format_options=format_options,
    )

    with trace_span("docling.convert_external_llm", file=doc_path.name, provider=provider, model=model):
        logger.info("Converting %s (%s/%s)", doc_path.name, provider, model)
        result = converter.convert(str(doc_path))
        logger.info("Status: %s", result.status)

    return result.document.export_to_markdown()  # type: ignore[no-any-return]


def convert_docx(  # noqa: PLR0912, PLR0915
    docx_path: Path,
    output_dir: Path,
    *,
    vlm: bool,
    vlm_preset: str,
    figures: bool,
    all_formats: bool,
    external_llm: tuple[str, str] | None,
    settings: Settings,
) -> None:
    """DOCX conversion: Docling native + image extraction + VLM descriptions."""
    output_dir.mkdir(parents=True, exist_ok=True)

    converter = DocumentConverter(
        allowed_formats=[InputFormat.DOCX],
        format_options={InputFormat.DOCX: WordFormatOption()},
    )

    with trace_span("docling.convert_docx", file=docx_path.name):
        logger.info("Converting %s (DOCX native)", docx_path.name)
        result = converter.convert(str(docx_path))
        logger.info("Status: %s", result.status)

    doc = result.document

    figure_map: dict[str, str] = {}
    figure_descriptions: dict[str, str] = {}
    fig_count = 0
    image_paths: list[Path] = []
    item_refs: list[str] = []

    if figures:
        fig_dir = output_dir / "figures"
        fig_dir.mkdir(exist_ok=True)
        for item, _ in doc.iterate_items():
            if isinstance(item, PictureItem):
                img = item.get_image(doc)
                if img:
                    filename = f"figure_{fig_count}.png"
                    img.save(fig_dir / filename)
                    figure_map[item.self_ref] = f"figures/{filename}"
                    image_paths.append(fig_dir / filename)
                    item_refs.append(item.self_ref)
                    fig_count += 1

        if vlm and image_paths:
            with trace_span("vlm.describe_docx_images", count=len(image_paths)):
                if external_llm:
                    provider, model = external_llm
                    desc_list = describe_images_with_external_llm(image_paths, provider, model, settings)
                else:
                    desc_list = describe_images_with_vlm(image_paths, vlm_preset, Path(settings.models_path))
                for ref, desc in zip(item_refs, desc_list, strict=True):
                    if desc:
                        figure_descriptions[ref] = desc

    title = ""
    for t in doc.texts:
        if t.label == DocItemLabel.TITLE:
            title = t.text
            break

    page_md = build_page_annotated_markdown(doc, figure_map, title, figure_descriptions=figure_descriptions)
    (output_dir / "document.md").write_text(page_md)

    if fig_count > 0:
        catalog = build_images_catalog(doc, figure_map)
        (output_dir / "images.md").write_text(catalog)

    if all_formats:
        (output_dir / "output.md").write_text(doc.export_to_markdown())
        (output_dir / "output.json").write_text(json.dumps(doc.export_to_dict(), indent=2, default=str))
        (output_dir / "output.txt").write_text(doc.export_to_text())
        (output_dir / "output.html").write_text(doc.export_to_html())

    if title:
        logger.info("Title: %s", title)
    _print_output_summary(output_dir, fig_count, all_formats, vlm and fig_count > 0, len(figure_descriptions))


def convert_xlsx(xlsx_path: Path, output_dir: Path, *, all_formats: bool) -> None:
    """XLSX conversion: Docling native to directory output."""
    output_dir.mkdir(parents=True, exist_ok=True)

    converter = DocumentConverter(
        allowed_formats=[InputFormat.XLSX],
        format_options={InputFormat.XLSX: ExcelFormatOption()},
    )

    with trace_span("docling.convert_xlsx", file=xlsx_path.name):
        logger.info("Converting %s (XLSX native)", xlsx_path.name)
        result = converter.convert(str(xlsx_path))
        logger.info("Status: %s", result.status)

    doc = result.document
    (output_dir / "document.md").write_text(doc.export_to_markdown())

    if all_formats:
        (output_dir / "output.json").write_text(json.dumps(doc.export_to_dict(), indent=2, default=str))
        (output_dir / "output.txt").write_text(doc.export_to_text())
        (output_dir / "output.html").write_text(doc.export_to_html())

    _print_output_summary(output_dir, 0, all_formats)


def _map_pptx_images_to_items(
    doc: object,
    slide_images: dict[int, list[Path]],
) -> tuple[dict[str, str], list[Path], list[str]]:
    """Map python-pptx extracted images to Docling PictureItems by slide + order.

    Returns (figure_map, image_paths, item_refs) where:
      figure_map: {item.self_ref: relative_path}
      image_paths: ordered list of image file paths
      item_refs: ordered list of item self_refs (parallel to image_paths)
    """
    figure_map: dict[str, str] = {}
    image_paths: list[Path] = []
    item_refs: list[str] = []
    picture_idx_per_slide: dict[int, int] = {}

    for item, _ in doc.iterate_items():  # type: ignore[attr-defined]
        if isinstance(item, PictureItem):
            prov = getattr(item, "prov", None)
            slide_num = prov[0].page_no if prov else 0
            idx = picture_idx_per_slide.get(slide_num, 0)
            picture_idx_per_slide[slide_num] = idx + 1

            slide_imgs = slide_images.get(slide_num, [])
            if idx < len(slide_imgs):
                figure_map[item.self_ref] = f"figures/{slide_imgs[idx].name}"
                image_paths.append(slide_imgs[idx])
                item_refs.append(item.self_ref)

    return figure_map, image_paths, item_refs


def convert_pptx(  # noqa: PLR0912
    pptx_path: Path,
    output_dir: Path,
    *,
    vlm: bool,
    vlm_preset: str,
    figures: bool,
    all_formats: bool,
    external_llm: tuple[str, str] | None,
    settings: Settings,
) -> None:
    """PPTX conversion: Docling native text + python-pptx images + VLM descriptions."""
    output_dir.mkdir(parents=True, exist_ok=True)

    converter = DocumentConverter(
        allowed_formats=[InputFormat.PPTX],
        format_options={InputFormat.PPTX: PowerpointFormatOption()},
    )

    with trace_span("docling.convert_pptx", file=pptx_path.name):
        logger.info("Converting %s (PPTX native)", pptx_path.name)
        result = converter.convert(str(pptx_path))
        logger.info("Status: %s", result.status)

    doc = result.document

    figure_map: dict[str, str] = {}
    figure_descriptions: dict[str, str] = {}
    all_image_paths: list[Path] = []

    if figures:
        slide_images = extract_pptx_images(pptx_path, output_dir)
        figure_map, all_image_paths, item_refs = _map_pptx_images_to_items(doc, slide_images)

        if vlm and all_image_paths:
            with trace_span("vlm.describe_pptx_images", count=len(all_image_paths)):
                if external_llm:
                    provider, model = external_llm
                    desc_list = describe_images_with_external_llm(all_image_paths, provider, model, settings)
                else:
                    desc_list = describe_images_with_vlm(all_image_paths, vlm_preset, Path(settings.models_path))
                for ref, desc in zip(item_refs, desc_list, strict=True):
                    if desc:
                        figure_descriptions[ref] = desc

    title = ""
    for t in doc.texts:
        if t.label == DocItemLabel.TITLE:
            title = t.text
            break

    page_md = build_page_annotated_markdown(doc, figure_map, title, figure_descriptions=figure_descriptions)
    (output_dir / "document.md").write_text(page_md)

    fig_count = len(all_image_paths)
    if fig_count > 0:
        catalog = build_images_catalog(doc, figure_map)
        (output_dir / "images.md").write_text(catalog)

    if all_formats:
        (output_dir / "output.md").write_text(doc.export_to_markdown())
        (output_dir / "output.json").write_text(json.dumps(doc.export_to_dict(), indent=2, default=str))
        (output_dir / "output.txt").write_text(doc.export_to_text())
        (output_dir / "output.html").write_text(doc.export_to_html())

    if title:
        logger.info("Title: %s", title)
    console.print(f"[green]Output:[/green] {output_dir}/")
    console.print("  document.md  (slide-annotated markdown)")
    if fig_count > 0:
        console.print("  images.md    (image catalog)")
        console.print(f"  figures/     ({fig_count} figure(s))")
        if vlm:
            console.print(f"  VLM descriptions: {len(figure_descriptions)}/{fig_count}")
    if all_formats:
        console.print("  output.*     (md, html, json, txt)")


# ── Output helpers ────────────────────────────────────────────────────────────


def _print_output_summary(
    output_dir: Path,
    fig_count: int = 0,
    all_formats: bool = False,
    vlm_used: bool = False,
    desc_count: int = 0,
    extra_files: list[str] | None = None,
) -> None:
    """Print a consistent output summary."""
    console.print(f"[green]Output:[/green] {output_dir}/")
    console.print("  document.md")
    if fig_count > 0:
        console.print("  images.md")
        console.print(f"  figures/     ({fig_count} figure(s))")
        if vlm_used:
            console.print(f"  VLM descriptions: {desc_count}/{fig_count}")
    for f in extra_files or []:
        console.print(f"  {f}")
    if all_formats:
        console.print("  output.*     (md, html, json, txt)")


def _resolve_output_dir(source_path: Path | None, name: str, output_override: str | None) -> Path:
    """Compute the <name>_docling/ output directory."""
    if output_override:
        return Path(output_override)
    parent = source_path.parent if source_path else Path.cwd()
    return parent / f"{name}_docling"


# ── CLI ──────────────────────────────────────────────────────────────────────


def _download_models(preset: str, settings: Settings) -> None:
    """Download VLM model for offline use."""
    from docling.models.utils.hf_model_download import download_hf_model  # noqa: PLC0415

    dest = Path(settings.models_path)
    dest.mkdir(parents=True, exist_ok=True)

    repo_id = PRESET_REPO_IDS[preset]
    cache_folder = repo_id.replace("/", "--")
    target = dest / cache_folder

    if target.exists():
        console.print(f"[yellow]Already downloaded:[/yellow] {target}")
        return

    console.print(f"Downloading [bold]{repo_id}[/bold] to {dest}/")
    download_hf_model(repo_id=repo_id, local_dir=target, progress=True)
    console.print(f"[green]Done.[/green] Model available at {target}")


DEFAULT_MEDIA_LLM = "openrouter/google/gemini-2.5-flash"


def _resolve_media_llm(use_external_llm: str | None, settings: Settings) -> tuple[str, str, str]:
    """Resolve provider, model, and API key for audio/video processing."""
    llm_spec = use_external_llm or DEFAULT_MEDIA_LLM
    provider, model = parse_external_llm(llm_spec)
    api_key = _require_api_key(provider, settings)
    return provider, model, api_key


def _check_cache(out_path: Path, force: bool) -> bool:
    """Return True if output exists and force is False (should skip)."""
    if out_path.exists() and not force:
        has_content = any(out_path.iterdir()) if out_path.is_dir() else out_path.stat().st_size > 0
        if has_content:
            console.print(f"[yellow]Output already exists:[/yellow] {out_path}")
            console.print("[dim]Use -f to force re-conversion[/dim]")
            return True
    return False


@app.command()
def main(  # noqa: PLR0912, PLR0915
    document: str = typer.Argument(
        None,
        help="File path or URL. Required unless --start-audio or --download-models is used.",
    ),
    output: str | None = typer.Option(
        None,
        "-o",
        "--output",
        help="Override output directory (default: <name>_docling/)",
    ),
    use_external_llm: str | None = typer.Option(
        None,
        "--use-external-llm",
        help="External LLM provider: google/<model> or openrouter/<model>",
    ),
    vlm_preset: VlmPreset = typer.Option(
        VlmPreset.smolvlm,
        "--vlm-preset",
        help="Local VLM preset for picture descriptions",
    ),
    no_ocr: bool = typer.Option(
        False,
        "--no-ocr",
        help="Disable OCR (PDF only)",
    ),
    no_vlm: bool = typer.Option(
        False,
        "--no-vlm",
        help="Disable VLM picture descriptions (PDF, DOCX, PPTX)",
    ),
    no_figures: bool = typer.Option(
        False,
        "--no-figures",
        help="Skip figure extraction (text + tables only, faster)",
    ),
    all_formats: bool = typer.Option(
        False,
        "--all",
        help="Export all formats (md, html, json, txt) in addition to document.md",
    ),
    start_audio: bool = typer.Option(
        False,
        "--start-audio",
        help="Record from microphone (Ctrl+C to stop). DOCUMENT becomes the name.",
    ),
    analyze: bool = typer.Option(
        False,
        "--analyze",
        help="Add analysis pass: analysis.md (audio: summary, video: executive brief)",
    ),
    meeting: str | None = typer.Option(
        None,
        "-m",
        "--meeting",
        help="Meeting name or context for audio/video prompts",
    ),
    instructions: str | None = typer.Option(
        None,
        "-i",
        "--instructions",
        help="Custom prompt for --analyze (overrides default analysis prompt)",
    ),
    force: bool = typer.Option(
        False,
        "-f",
        "--force",
        help="Force re-conversion even if output already exists",
    ),
    download_models: bool = typer.Option(
        False,
        "--download-models",
        help="Download the VLM model selected by --vlm-preset for offline use",
    ),
    verbose: int = typer.Option(
        0,
        "-v",
        "--verbose",
        count=True,
        help="Increase verbosity (-vv for debug)",
    ),
    quiet: bool = typer.Option(
        False,
        "-q",
        "--quiet",
        help="Only show warnings and errors",
    ),
) -> None:
    """Convert documents, audio, and video to markdown.

    Output always goes to <name>_docling/ directory with document.md as main file.
    Use -o to override the output directory.

    \b
    Documents:  doc-convert document.pdf
    Audio:      doc-convert meeting.ogg [--analyze]
    Video:      doc-convert video.mp4 [--analyze]
    YouTube:    doc-convert https://youtube.com/watch?v=...
    Record:     doc-convert --start-audio "Meeting Name"
    Models:     doc-convert --download-models
    """
    setup_logging(verbose=verbose, quiet=quiet)

    if download_models:
        _download_models(vlm_preset.value, Settings())
        raise typer.Exit()

    # ── Audio recording mode ─────────────────────────────────────────────
    if start_audio:
        if not document:
            console.print('[red]--start-audio requires a name: doc-convert --start-audio "Meeting Name"[/red]')
            raise typer.Exit(1)

        from audio import record_audio  # noqa: PLC0415

        configure_tracing("doc-convert")
        settings = Settings()
        out_dir = _resolve_output_dir(None, document, output)
        out_dir.mkdir(parents=True, exist_ok=True)
        audio_file = record_audio(out_dir / "audio.ogg")
        _process_media_to_dir(
            audio_file, out_dir, "audio", meeting or document, analyze, instructions, force, use_external_llm, settings
        )
        raise typer.Exit()

    if document is None:
        console.print("Missing argument 'DOCUMENT'. See --help.")
        raise typer.Exit(1)

    configure_tracing("doc-convert")
    settings = Settings()
    models = Path(settings.models_path)
    ext_llm = parse_external_llm(use_external_llm) if use_external_llm else None

    # ── YouTube URL ──────────────────────────────────────────────────────
    from video import is_youtube_url  # noqa: PLC0415

    if is_youtube_url(document):
        from video import download_youtube  # noqa: PLC0415

        downloaded = download_youtube(document)
        out_dir = _resolve_output_dir(Path.cwd(), downloaded.stem, output)
        try:
            _process_media_to_dir(
                downloaded, out_dir, "video", meeting, analyze, instructions, force, use_external_llm, settings
            )
        finally:
            downloaded.unlink(missing_ok=True)
        raise typer.Exit()

    # ── File-based conversion ────────────────────────────────────────────
    tmp_file: Path | None = None
    try:
        if is_google_url(document):
            tmp_file, title, fmt = download_google_doc(document, settings)
            doc_path = tmp_file
            doc_name = title
        else:
            doc_path = Path(document)
            if not doc_path.exists():
                console.print(f"[red]{doc_path} not found[/red]")
                raise typer.Exit(1)

            from media_llm import is_audio_ext, is_video_ext  # noqa: PLC0415

            ext = doc_path.suffix.lower()
            if is_audio_ext(ext):
                out_dir = _resolve_output_dir(doc_path, doc_path.stem, output)
                _process_media_to_dir(
                    doc_path, out_dir, "audio", meeting, analyze, instructions, force, use_external_llm, settings
                )
                raise typer.Exit()
            if is_video_ext(ext):
                out_dir = _resolve_output_dir(doc_path, doc_path.stem, output)
                _process_media_to_dir(
                    doc_path, out_dir, "video", meeting, analyze, instructions, force, use_external_llm, settings
                )
                raise typer.Exit()

            fmt = detect_format(doc_path)
            doc_name = doc_path.stem

        out_dir = _resolve_output_dir(doc_path, doc_name, output)
        if _check_cache(out_dir, force):
            raise typer.Exit()

        # ── Document conversion (Docling) ────────────────────────────────
        if fmt == InputFormat.PDF and not ext_llm:
            convert_pdf_local(
                pdf_path=doc_path,
                output_dir=out_dir,
                do_ocr=not no_ocr,
                vlm=not no_vlm,
                vlm_preset=vlm_preset.value,
                figures=not no_figures,
                all_formats=all_formats,
                models_path=models,
            )
        elif fmt == InputFormat.PPTX:
            convert_pptx(
                pptx_path=doc_path,
                output_dir=out_dir,
                vlm=not no_vlm,
                vlm_preset=vlm_preset.value,
                figures=not no_figures,
                all_formats=all_formats,
                external_llm=ext_llm,
                settings=settings,
            )
        elif fmt == InputFormat.DOCX:
            convert_docx(
                docx_path=doc_path,
                output_dir=out_dir,
                vlm=not no_vlm,
                vlm_preset=vlm_preset.value,
                figures=not no_figures,
                all_formats=all_formats,
                external_llm=ext_llm,
                settings=settings,
            )
        elif fmt == InputFormat.XLSX:
            convert_xlsx(xlsx_path=doc_path, output_dir=out_dir, all_formats=all_formats)
        elif fmt == InputFormat.IMAGE:
            if not ext_llm:
                console.print("[red]Image conversion requires --use-external-llm[/red]")
                raise typer.Exit(1)
            provider, model = ext_llm
            md = convert_with_external_llm(doc_path, fmt, provider, model, settings)
            out_dir.mkdir(parents=True, exist_ok=True)
            (out_dir / "document.md").write_text(md)
            _print_output_summary(out_dir)
        elif fmt == InputFormat.PDF and ext_llm:
            provider, model = ext_llm
            md = convert_with_external_llm(doc_path, fmt, provider, model, settings)
            out_dir.mkdir(parents=True, exist_ok=True)
            (out_dir / "document.md").write_text(md)
            _print_output_summary(out_dir)
        else:
            console.print(f"[red]Unsupported format: {fmt}[/red]")
            raise typer.Exit(1)
    finally:
        if tmp_file and tmp_file.exists():
            tmp_file.unlink()


# ── Audio/Video processing ───────────────────────────────────────────────────


def _process_media_to_dir(
    media_path: Path,
    output_dir: Path,
    media_type: str,
    meeting: str | None,
    analyze: bool,
    instructions: str | None,
    force: bool,
    use_external_llm: str | None,
    settings: Settings,
) -> None:
    """Process audio or video into a _docling/ directory with document.md + optional analysis.md."""
    from media_llm import process_media  # noqa: PLC0415

    if _check_cache(output_dir, force):
        return

    output_dir.mkdir(parents=True, exist_ok=True)
    provider, model, api_key = _resolve_media_llm(use_external_llm, settings)

    # Main conversion (transcription for audio, extraction for video)
    if media_type == "audio":
        from audio import build_transcription_prompt  # noqa: PLC0415

        prompt, system = build_transcription_prompt(meeting)
        span_name = "audio.transcribe"
    else:
        from video import build_extraction_prompt  # noqa: PLC0415

        prompt, system = build_extraction_prompt(meeting)
        span_name = "video.extract"

    with trace_span(span_name, file=media_path.name, provider=provider):
        md = process_media(media_path, provider, model, prompt, api_key, system_prompt=system)
    if meeting:
        md = f"# {meeting}\n\n{md}"
    (output_dir / "document.md").write_text(md)

    extra: list[str] = []

    # Analysis (optional second pass)
    if analyze:
        if instructions:
            a_prompt = instructions
            a_system = f"Context: {meeting}\n\n{instructions}" if meeting else instructions
        elif media_type == "audio":
            from audio import build_analysis_prompt as audio_analysis  # noqa: PLC0415

            a_prompt, a_system = audio_analysis(meeting)
        else:
            from video import build_analysis_prompt as video_analysis  # noqa: PLC0415

            a_prompt, a_system = video_analysis(meeting)

        with trace_span(f"{media_type}.analyze", file=media_path.name, provider=provider):
            analysis_md = process_media(media_path, provider, model, a_prompt, api_key, system_prompt=a_system)
        if meeting:
            analysis_md = f"# {meeting}\n\n{analysis_md}"
        (output_dir / "analysis.md").write_text(analysis_md)
        extra.append("analysis.md")

    _print_output_summary(output_dir, extra_files=extra)
